"""Document processor service for extracting text from PowerPoint files."""

import os
import logging
from typing import Dict, List, Optional, Any, Tuple
import tempfile
import io

import pandas as pd
import numpy as np
from pptx import Presentation
from googleapiclient.http import MediaIoBaseDownload

from app.config import settings
from app.services.auth import get_google_drive_service

logger = logging.getLogger(__name__)

class DocumentProcessor:
    """Service for processing PowerPoint documents from Google Drive."""
    
    def __init__(self):
        """Initialize the document processor."""
        self.drive_service = None
        self.data_dir = settings.DATA_DIR
        
        # Create data directory if it doesn't exist
        os.makedirs(self.data_dir, exist_ok=True)
    
    def _init_drive_service(self):
        """Initialize the Google Drive service if not already initialized."""
        if not self.drive_service:
            self.drive_service = get_google_drive_service()
    
    async def list_drive_files(self, folder_id: Optional[str] = None) -> List[Dict]:
        """List files in the specified Google Drive folder."""
        self._init_drive_service()
        
        if not folder_id:
            folder_id = settings.GOOGLE_DRIVE_FOLDER_ID
        
        try:
            query = f"'{folder_id}' in parents and trashed = false"
            fields = "files(id, name, mimeType, createdTime, modifiedTime)"
            
            response = self.drive_service.files().list(
                q=query,
                spaces='drive',
                fields=fields
            ).execute()
            
            return response.get('files', [])
        except Exception as e:
            logger.error(f"Error listing Drive files: {str(e)}")
            raise
    
    async def download_file(self, file_id: str) -> Tuple[str, bytes]:
        """Download a file from Google Drive."""
        self._init_drive_service()
        
        try:
            # Get file metadata
            file_metadata = self.drive_service.files().get(fileId=file_id).execute()
            file_name = file_metadata['name']
            
            # Download file content
            request = self.drive_service.files().get_media(fileId=file_id)
            file_content = io.BytesIO()
            downloader = MediaIoBaseDownload(file_content, request)
            
            done = False
            while not done:
                status, done = downloader.next_chunk()
            
            return file_name, file_content.getvalue()
        except Exception as e:
            logger.error(f"Error downloading file {file_id}: {str(e)}")
            raise
    
    async def process_ppt_file(self, file_content: bytes) -> str:
        """Extract text content from a PowerPoint file."""
        try:
            # Create a temporary file to save the PPT content
            with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as temp_file:
                temp_file.write(file_content)
                temp_file_path = temp_file.name
            
            # Extract text from the PPT file
            text_content = self._extract_text_from_ppt(temp_file_path)
            
            # Clean up the temporary file
            os.unlink(temp_file_path)
            
            return text_content
        except Exception as e:
            logger.error(f"Error processing PPT file: {str(e)}")
            raise
    
    def _extract_text_from_ppt(self, file_path: str) -> str:
        """Extract text from a PowerPoint file."""
        presentation = Presentation(file_path)
        text_content = []
        
        # Extract slide titles and content
        for i, slide in enumerate(presentation.slides):
            slide_text = []
            
            # Get slide title if available
            if slide.shapes.title and slide.shapes.title.text:
                slide_text.append(f"Slide {i+1} - {slide.shapes.title.text}")
            else:
                slide_text.append(f"Slide {i+1}")
            
            # Extract text from all shapes in the slide
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    # Skip if this is the title we already added
                    if shape == slide.shapes.title:
                        continue
                    slide_text.append(shape.text)
            
            # Add slide content to overall content
            text_content.append("\n".join(slide_text))
        
        return "\n\n".join(text_content)
    
    async def process_all_files(self, folder_id: Optional[str] = None) -> Dict[str, str]:
        """Process all PowerPoint files in the specified folder."""
        # List all files in the folder
        files = await self.list_drive_files(folder_id)
        
        # Filter for PowerPoint files
        ppt_files = [f for f in files if f['mimeType'] == 'application/vnd.google-apps.presentation' or 
                    f['mimeType'] == 'application/vnd.openxmlformats-officedocument.presentationml.presentation']
        
        results = {}
        for file in ppt_files:
            try:
                # Download the file
                file_name, file_content = await self.download_file(file['id'])
                
                # Process the file
                text_content = await self.process_ppt_file(file_content)
                
                # Store the result
                results[file_name] = {
                    'id': file['id'],
                    'name': file_name,
                    'content': text_content,
                    'metadata': {
                        'created': file.get('createdTime'),
                        'modified': file.get('modifiedTime')
                    }
                }
                
                logger.info(f"Successfully processed file: {file_name}")
            except Exception as e:
                logger.error(f"Error processing file {file.get('name')}: {str(e)}")
        
        return results